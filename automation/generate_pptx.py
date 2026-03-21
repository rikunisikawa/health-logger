#!/usr/bin/env python3
"""
automation/generate_pptx.py
スライド原稿 JSON から PowerPoint ファイルを生成する。

使い方:
  python3 automation/generate_pptx.py docs/explain/20260321-120000.json
  python3 automation/generate_pptx.py docs/explain/20260321-120000.json --output docs/explain/report.pptx

依存:
  pip install python-pptx
"""

import argparse
import json
import sys
from pathlib import Path


def generate(json_path: Path, output_path: Path | None = None) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("ERROR: python-pptx が必要です: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    deck = json.loads(json_path.read_text(encoding="utf-8"))

    if output_path is None:
        output_path = json_path.with_suffix(".pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
    ACCENT_COLOR = RGBColor(0x16, 0x21, 0x3E)
    TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
    BULLET_COLOR = RGBColor(0x22, 0x5A, 0x84)

    blank_layout = prs.slide_layouts[6]  # 空白レイアウト

    def add_textbox(slide, text, left, top, width, height,
                    font_size=18, bold=False, color=TEXT_COLOR, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    for slide_data in deck.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # 背景（タイトルバー）
        from pptx.util import Emu
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(1.4)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = TITLE_COLOR
        bg.line.fill.background()

        # スライドタイトル
        add_textbox(
            slide, slide_data.get("title", ""),
            left=0.3, top=0.2, width=12.5, height=1.0,
            font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.LEFT,
        )

        # 箇条書き
        bullets = slide_data.get("bullets", [])
        y = 1.7
        for bullet in bullets:
            if not bullet:
                continue
            add_textbox(
                slide, f"• {bullet}",
                left=0.5, top=y, width=11.8, height=0.6,
                font_size=18, color=BULLET_COLOR,
            )
            y += 0.65
            if y > 6.2:
                break

        # ナレーション（フッター）
        narration = slide_data.get("narration", "")
        if narration:
            narration_box = slide.shapes.add_shape(
                1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
            narration_box.fill.solid()
            narration_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            narration_box.line.fill.background()
            add_textbox(
                slide, f"📢 {narration}",
                left=0.3, top=6.85, width=12.7, height=0.5,
                font_size=11, color=RGBColor(0x66, 0x66, 0x66),
            )

        # スライド番号
        add_textbox(
            slide,
            f"{slide_data.get('slide_number', '')} / {len(deck['slides'])}",
            left=11.8, top=0.3, width=1.3, height=0.5,
            font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC),
            align=PP_ALIGN.RIGHT,
        )

    prs.save(str(output_path))
    print(f"✅ PowerPoint を生成しました: {output_path}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="スライド原稿 JSON から PowerPoint を生成")
    parser.add_argument("json_path", help="スライド原稿 JSON のパス")
    parser.add_argument("--output", help="出力先 pptx パス（省略時は JSON と同ディレクトリ）")
    args = parser.parse_args()

    json_path = Path(args.json_path)
    if not json_path.exists():
        print(f"ERROR: ファイルが見つかりません: {json_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output) if args.output else None
    generate(json_path, output_path)
